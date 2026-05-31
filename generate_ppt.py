import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Try to import PIL for aspect ratio detection
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("PIL/Pillow is not installed. Falling back to default scaling.")

# Define project working directory
proj_dir = "/Users/rennuoya/Desktop/RPA/Chinese Royals Fun"

def add_header(slide, title_text):
    # Theme Colors
    TEAL = RGBColor(22, 115, 108)
    GOLD = RGBColor(224, 182, 85)
    
    # 1. Left decorative bar (Gold)
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.6)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = GOLD
    left_bar.line.fill.background() # no border
    
    # 2. Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "SimHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    # 3. Horizontal subtle divider line
    div_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.15), Inches(12.133), Inches(0.02)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(226, 232, 240) # light gray border
    div_line.line.fill.background()

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_image_fit(slide, img_name, left, top, max_width, max_height):
    """Adds an image keeping its aspect ratio and centering it in the bounding box."""
    img_path = os.path.join(proj_dir, "assets/images", img_name)
    if not os.path.exists(img_path):
        print(f"Warning: Image file not found: {img_path}")
        return None
        
    if HAS_PIL:
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
            
            # Target aspect ratio
            target_aspect = max_width / max_height
            
            if aspect > target_aspect:
                # Image is wider than bounding box aspect
                width = max_width
                height = max_width / aspect
            else:
                # Image is taller than bounding box aspect
                height = max_height
                width = max_height * aspect
                
            # Center coordinates inside bounding box
            x_offset = (max_width - width) / 2
            y_offset = (max_height - height) / 2
            
            return slide.shapes.add_picture(
                img_path, 
                left + x_offset, 
                top + y_offset, 
                width=width, 
                height=height
            )
        except Exception as e:
            print(f"PIL processing error for {img_name}: {e}. Falling back to default scaling.")
            
    # Default fallback
    try:
        return slide.shapes.add_picture(img_path, left, top, width=max_width)
    except Exception as e:
        print(f"Failed to add image {img_name}: {e}")
    return None

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Common layout: blank
    blank_layout = prs.slide_layouts[6]
    
    # Premium Theme Colors (Option A: Teal & Gold)
    TEAL = RGBColor(22, 115, 108)      # Gemstone Teal
    CREAM = RGBColor(250, 248, 242)    # Warm Oat Sand Background
    GOLD = RGBColor(224, 182, 85)      # Champagne Gold
    WHITE = RGBColor(255, 255, 255)
    TEXT_MAIN = RGBColor(30, 41, 59)   # Slate Dark Gray
    TEXT_MUTED = RGBColor(71, 85, 105) # Slate Soft Gray
    
    # ==========================================
    # SLIDE 1: COVER SLIDE (Teal Background)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, TEAL)
    
    # Top Gold accent line
    top_accent = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.15)
    )
    top_accent.fill.solid()
    top_accent.fill.fore_color.rgb = GOLD
    top_accent.line.fill.background()
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "御苑福膳"
    p1.font.name = "SimHei"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = GOLD
    p1.alignment = PP_ALIGN.LEFT
    p1.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "专注本土超级食物 · 开启健康膳食新风尚"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT
    
    # Subtitle / Company Box
    company_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
    tf2 = company_box.text_frame
    tf2.word_wrap = True
    
    pc1 = tf2.paragraphs[0]
    pc1.text = "北京中卫神州生物科技有限公司 创立"
    pc1.font.name = "Microsoft YaHei"
    pc1.font.size = Pt(16)
    pc1.font.color.rgb = RGBColor(210, 230, 228)
    pc1.space_after = Pt(8)
    
    pc2 = tf2.add_paragraph()
    pc2.text = "北京御苑福膳文化传播有限公司 运营管理"
    pc2.font.name = "Microsoft YaHei"
    pc2.font.size = Pt(16)
    pc2.font.color.rgb = RGBColor(210, 230, 228)

    # ==========================================
    # SLIDE 2: BRAND ORIGIN & RESEARCH BACKGROUND (Cream Background)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, CREAM)
    add_header(slide2, "品牌起源与中医药科研背景")
    
    # Left Box - Description Card
    left_card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    left_tf = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.4)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "深厚中医药学科研背景"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(12)
    
    p = left_tf.add_paragraph()
    p.text = "“御苑福膳”由北京中卫神州生物科技有限公司（成立于2002年，原隶属中国医学科学院药用植物研究所企业）创立。现由其全资子公司北京御苑福膳文化传播有限公司进行品牌运营。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = left_tf.add_paragraph()
    p.text = "依托科研历史积淀，跨界联手营养学界与国家级烹饪大师，倾力构建餐桌健康工程，为公众的体质调理与健康饮食提供切实有效的膳食方案。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MAIN
    
    # Left Card Image (Research Base)
    add_image_fit(slide2, "base_research.jpg", Inches(0.8), Inches(4.0), Inches(5.4), Inches(2.4))
    
    # Right Box - Expert Guidance List
    right_card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    
    right_tf = slide2.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.4)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "国医大师与权威专家学术指导"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(12)
    
    bullet_points = [
        "中药学泰斗、工程院院士肖培根长期担纲学术指导专家，确保底蕴；",
        "国医大师、工程院院士王琦亲临指导，对时令膳食开发提供科学支撑；",
        "汇聚多位资深中医药学者、营养师与烹饪大师构建科学健康的膳食模式。"
    ]
    for bp in bullet_points:
        p = right_tf.add_paragraph()
        p.text = "• " + bp
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    # Right Card Images (Xiao Peigen & Wang Qi)
    add_image_fit(slide2, "xiao_peigen.jpg", Inches(7.1), Inches(3.9), Inches(2.5), Inches(2.1))
    add_image_fit(slide2, "wang_qi.jpg", Inches(9.9), Inches(3.9), Inches(2.5), Inches(2.1))
    
    # Right Card Image Captions
    cap_xiao = slide2.shapes.add_textbox(Inches(7.1), Inches(6.05), Inches(2.5), Inches(0.4)).text_frame
    cap_xiao.word_wrap = True
    p = cap_xiao.paragraphs[0]
    p.text = "肖培根 院士"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    cap_wang = slide2.shapes.add_textbox(Inches(9.9), Inches(6.05), Inches(2.5), Inches(0.4)).text_frame
    cap_wang.word_wrap = True
    p = cap_wang.paragraphs[0]
    p.text = "王琦 院士"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3: DIETARY CONCEPT (Cream Background)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, CREAM)
    add_header(slide3, "核心理念：本土超级食物与新鲜烹饪")
    
    # Left Column
    left_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    left_tf = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(5.0)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "新鲜食材 · 新新鲜烹饪"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(16)
    
    points_left = [
        "坚持本真膳食体验，拒绝工业预制与死板调味，完全遵循“新鲜食材、现场即时烹饪”的工艺准则；",
        "精选并推广高营养素密度的中国本土超级食物，将日常朴素食材以极致的烹饪艺术做出千姿百态的花样；",
        "紧密围绕自然四季、二十四节气的时令更替进行菜谱轮换，牢牢锁住新鲜食材的原始生机、活性成分与充沛营养；",
        "以健康为底色，把大众餐桌的日常饮食升级为兼具美味享受与健康调理的精致艺术体验。"
    ]
    for pt in points_left:
        p = left_tf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(14)
        
    # Right Column
    right_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    
    right_tf = slide3.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "中药服食与时令食疗创新"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(12)
    
    points_right = [
        "首创“鲜药入膳”核心专利模式，探索生鲜本草与现代烹饪的结合；",
        "直接依托高科技鲜中药材大棚，现采现用藿香、紫苏、玉竹等时令鲜药；",
        "科学融合中医药饮食理论与现代食疗临床营养成果，研发四季调理养生方。"
    ]
    for pt in points_right:
        p = right_tf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)
        
    # Right Card Image (Greenhouse Base)
    add_image_fit(slide3, "base_greenhouse.jpg", Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.2))
    
    # Image Caption
    cap_gh = slide3.shapes.add_textbox(Inches(7.0), Inches(6.22), Inches(5.5), Inches(0.4)).text_frame
    cap_gh.word_wrap = True
    p = cap_gh.paragraphs[0]
    p.text = "生态示范基地 · 现代化高科技鲜药材大棚"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 4: ECOLOGICAL BASES (Cream Background)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, CREAM)
    add_header(slide4, "生态种植基地与官方认定示范窗口")
    
    # Text card at top
    top_card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.9)
    )
    top_card.fill.solid()
    top_card.fill.fore_color.rgb = WHITE
    top_card.line.color.rgb = RGBColor(226, 232, 240)
    
    top_tf = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(1.7)).text_frame
    top_tf.word_wrap = True
    
    p = top_tf.paragraphs[0]
    p.text = "中医药文化旅游示范基地（北京市中医管理局官方认定）"
    p.font.name = "SimHei"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(8)
    
    p = top_tf.add_paragraph()
    p.text = "由北京市中医管理局与北京市旅游局联合评定、授权颁发官方示范牌照。它是经政府相关部门严格评审与行业权威认定的中医药养生与膳食科普示范标杆窗口，融植物科普、中药识药与健康膳食体验于一体。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MAIN
    
    # Top Card Image (Demo Garden Plaque)
    add_image_fit(slide4, "base_demogarden.jpg", Inches(8.5), Inches(1.5), Inches(4.0), Inches(1.7))
    
    # 4 Key Bases Grid at bottom (Image + Text)
    grid_start_y = Inches(3.5)
    grid_w = Inches(2.9)
    grid_h = Inches(3.4)
    gap_x = Inches(0.16)
    
    bases_subset = [
        ("原生态灵芝基地", "base_lingzhi.jpg", "深山林区环境抚育，活体孢子切片采收，饱含灵芝活性多糖营养。"),
        ("宁夏枸杞基地", "base_gouqi.jpg", "黄河冲积平原种植，肉厚汁甜，富含类胡萝卜素，清晨采摘鲜用。"),
        ("阿拉善肉苁蓉基地", "base_roucongrong.jpg", "沙漠边缘纯天然梭梭寄生，手工鲜取，提供极高活性与滋补动力。"),
        ("青藏玛咖基地", "base_maka.jpg", "高海拔强日照极寒山区抚育，汲取雪域精华，能量充沛，天然滋养。")
    ]
    
    for i, (title, img_name, desc) in enumerate(bases_subset):
        x = Inches(0.6) + i * (grid_w + gap_x)
        
        # Card Background
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, grid_start_y, grid_w, grid_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Card Image Box
        add_image_fit(slide4, img_name, x + Inches(0.1), grid_start_y + Inches(0.1), Inches(2.7), Inches(1.6))
        
        # Card Text Box
        tf = slide4.shapes.add_textbox(x + Inches(0.1), grid_start_y + Inches(1.75), Inches(2.7), Inches(1.5)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "SimHei"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 5: HONORS & QUALIFICATIONS (Cream Background)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, CREAM)
    add_header(slide5, "品牌资质与社会荣誉认可")
    
    # 5 Honors layout - 3 on top, 2 centered bottom
    honor_w = Inches(3.8)
    honor_h = Inches(2.6)
    gap_x = Inches(0.3)
    
    honors_list = [
        ("非物质文化遗产", "ich_certificate.jpg", "传统膳食技艺非遗传承保护单位", "多年来始终坚持中华古老膳食技艺的精髓继承与现代保护，非物质文化遗产的传承，见证了御苑福膳深厚的文化底蕴。"),
        ("老字号荣誉", "time_honored_brand.jpg", "传统膳食品牌与老字号荣誉", "荣获中医药及健康餐饮行业内的老字号名号，代表着数十载持之以恒的卓越声誉、工匠精神以及官方与社会的信任。"),
        ("行业荣誉资质", "honor_plaque.jpg", "特色康养餐饮与行业认证", "被行业协会与健康养生协会联合认证为药食同源与特色健康餐饮示范店，体现出深厚的草本膳食科学研究底蕴。"),
        ("首都窗口行业示范单位", "skills_demo.png", "2008年北京迎奥运技能示范", "在北京市劳动和社会保障局、首都窗口行业奥运工作协调小组联合评比中，被授予“首都窗口行业技能示范单位”牌匾。"),
        ("安全生产标准化企业", "safety_production.png", "餐饮业安全生产标准化企业", "作为北京市海淀区餐饮业安全生产标准化的典范，连续多年荣获安全生产达标资质，用极高安全标准保障每一位食客的健康安全。")
    ]
    
    for i, (title, img_name, subtitle, desc) in enumerate(honors_list):
        if i < 3:
            x = Inches(0.6) + i * (honor_w + gap_x)
            y = Inches(1.4)
        else:
            # Centered on bottom row
            x = Inches(2.65) + (i - 3) * (honor_w + gap_x)
            y = Inches(4.2)
            
        # Card background
        card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, honor_w, honor_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Embedded certificate/plaque image (vertical frame aspect)
        add_image_fit(slide5, img_name, x + Inches(0.15), y + Inches(0.25), Inches(1.3), Inches(2.1))
        
        # Text Frame
        tf = slide5.shapes.add_textbox(x + Inches(1.5), y + Inches(0.15), Inches(2.2), Inches(2.3)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "SimHei"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: CELEBRITY FOOTPRINTS (Cream Background)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, CREAM)
    add_header(slide6, "社会各界与中医药学专家交流足迹")
    
    # Left Box - Political leaders
    left_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    left_tf = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.6)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "党政领导与社会名流莅临视察"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(12)
    
    leaders = [
        "全国人大常委会副委员长 何维 莅临视察指导工作；",
        "全国政协原副主席 张思卿 & 孙孚凌 视察指导工作；",
        "原卫生部部长 张文康 & 高强、原副部长 顾英奇 亲临指导；",
        "甘肃省委原书记 陆浩、海南省委原书记 卫留成 等先后莅临；",
        "毛泽东主席原机要秘书 张玉凤 女士及刘爱民先生亲临品鉴。"
    ]
    for ld in leaders:
        p = left_tf.add_paragraph()
        p.text = "★ " + ld
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(5)
        
    # Left Column Photos (He Wei & Lu Hao)
    add_image_fit(slide6, "he_wei.jpg", Inches(0.8), Inches(4.2), Inches(2.6), Inches(2.0))
    add_image_fit(slide6, "lu_hao.jpg", Inches(3.6), Inches(4.2), Inches(2.6), Inches(2.0))
    
    cap_he = slide6.shapes.add_textbox(Inches(0.8), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_he.paragraphs[0]
    p.text = "何维 副委员长莅临考察"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    cap_lu = slide6.shapes.add_textbox(Inches(3.6), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_lu.paragraphs[0]
    p.text = "原省委书记 陆浩视察留影"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
        
    # Right Box - Experts & Masters
    right_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    
    right_tf = slide6.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.6)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "学术泰斗与国宴大师技术合作"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(12)
    
    experts = [
        "中药学权威 肖培根院士 亲自指导，从植物药理确保功效；",
        "国医大师 王琦院士 等泰斗级专家对药食同源应用予以厚赞；",
        "医科院原院长 顾方舟 & 刘德培院士 等多位科学家高度赞誉；",
        "国宴泰斗 侯仲华 & 孟宪斌大师 指导，融入最高级别烹饪技艺；",
        "接待联合国北北合作组织、西藏活佛及多国使节开展文化交流。"
    ]
    for exp in experts:
        p = right_tf.add_paragraph()
        p.text = "★ " + exp
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(5)
        
    # Right Column Photos (Liu Depei & Hou Zhonghua)
    add_image_fit(slide6, "liu_depei.jpg", Inches(7.0), Inches(4.2), Inches(2.6), Inches(2.0))
    add_image_fit(slide6, "hou_zhonghua.jpg", Inches(9.8), Inches(4.2), Inches(2.6), Inches(2.0))
    
    cap_liu = slide6.shapes.add_textbox(Inches(7.0), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_liu.paragraphs[0]
    p.text = "刘德培 院士亲临指导"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    cap_hou = slide6.shapes.add_textbox(Inches(9.8), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_hou.paragraphs[0]
    p.text = "国宴烹饪大师 侯仲华指导"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 7: SOCIAL RESPONSIBILITY & MEDIA (Cream Background)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, CREAM)
    add_header(slide7, "大健康政策响应、社会责任与媒体聚焦")
    
    # Left Column
    left_card = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(226, 232, 240)
    
    left_tf = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.6)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "积极响应政策与践行社会责任"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(10)
    
    p = left_tf.add_paragraph()
    p.text = "「大健康战略与体重管理年」：开设时令食育工坊，在端午、中秋、清明、儿童节等节日，联动社区普及低卡低脂少盐的营养餐，践行大健康理念。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(8)
    
    p = left_tf.add_paragraph()
    p.text = "「“爱心拐杖”云端老年食堂」：首批定点保障单位，专为海淀马连洼长者研发控盐控脂、软烂消化的温情膳食，守护银发族餐桌健康。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MAIN
    
    # Left Card Image (Academic Conference Group Photo)
    add_image_fit(slide7, "academic_conference.jpg", Inches(1.1), Inches(4.2), Inches(4.8), Inches(2.0))
    
    cap_conf = slide7.shapes.add_textbox(Inches(1.1), Inches(6.22), Inches(4.8), Inches(0.45)).text_frame
    p = cap_conf.paragraphs[0]
    p.text = "中医药膳学术研讨与品鉴会现场合影"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Right Column - Media Coverage
    right_card = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(226, 232, 240)
    
    right_tf = slide7.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.6)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "主流媒体聚焦与社会公信力"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(10)
    
    media = [
        ("CCTV-1《夕阳红》清明/谷雨/立夏/小满美食特别节目多次专题报道；", ""),
        ("CGTN西语频道向全球播放《民生中国：药食同源的传统新实践》；", ""),
        ("CCTV央视频 开展十一假日健康药食同源美食现场直播制作；", ""),
        ("北京卫视《气象观天下》连续录制端午缠角粽、惊蛰龙抬头节气美食。"
        "", "")
    ]
    for name, _ in media:
        if name:
            p = right_tf.add_paragraph()
            p.text = "● " + name
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(5)
            
    # Right Card Image (BTV Host He Beiqi Photo)
    add_image_fit(slide7, "he_beiqi.jpg", Inches(7.3), Inches(4.2), Inches(4.8), Inches(2.0))
    
    cap_bq = slide7.shapes.add_textbox(Inches(7.3), Inches(6.22), Inches(4.8), Inches(0.45)).text_frame
    p = cap_bq.paragraphs[0]
    p.text = "北京电视台主持人贺贝奇录制现场"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 8: END SLIDE (Teal Background)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, TEAL)
    
    # Bottom Gold accent line
    bottom_accent = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35), Inches(13.333), Inches(0.15)
    )
    bottom_accent.fill.solid()
    bottom_accent.fill.fore_color.rgb = GOLD
    bottom_accent.line.fill.background()
    
    # Text Box
    end_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf_end = end_box.text_frame
    tf_end.word_wrap = True
    
    p = tf_end.paragraphs[0]
    p.text = "坚持新鲜食材与新鲜烹饪"
    p.font.name = "SimHei"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)
    
    p = tf_end.add_paragraph()
    p.text = "以时令膳食守护大众餐桌健康"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(40)
    
    p = tf_end.add_paragraph()
    p.text = "谢谢观看"
    p.font.name = "SimHei"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filepath = os.path.join(proj_dir, "御苑福膳_品牌介绍.pptx")
    prs.save(filepath)
    print(f"Presentation generated successfully at {filepath}")

if __name__ == "__main__":
    create_presentation()
