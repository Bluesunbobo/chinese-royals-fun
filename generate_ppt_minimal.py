import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Try to import PIL for aspect ratio detection
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("PIL/Pillow is not installed. Falling back to default scaling.")

# Define project working directory
proj_dir = "/Users/rennuoya/Desktop/RPA/Chinese Royals Fun"

def add_header(slide, title_text):
    # Theme Colors (Option B: Terracotta & Sage)
    TERRACOTTA = RGBColor(186, 85, 66)
    SAGE = RGBColor(110, 145, 115)
    
    # 1. Left decorative bar (Sage Green)
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.6)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = SAGE
    left_bar.line.fill.background() # no border
    
    # 2. Slide Title (Terracotta)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "SimHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    
    # 3. Horizontal subtle divider line
    div_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.15), Inches(12.133), Inches(0.02)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = RGBColor(230, 226, 218) # warm beige border
    div_line.line.fill.background()

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_image_fit(slide, img_name, left, top, max_width, max_height):
    """Adds an image keeping its aspect ratio and centering it in the bounding box."""
    img_path = os.path.join(proj_dir, "assets/images", img_name)
    if not os.path.exists(img_path):
        print(f"Warning: Image file not found: {img_path}")
        return None
        
    if HAS_PIL:
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
            
            # Target aspect ratio
            target_aspect = max_width / max_height
            
            if aspect > target_aspect:
                # Image is wider than bounding box aspect
                width = max_width
                height = max_width / aspect
            else:
                # Image is taller than bounding box aspect
                height = max_height
                width = max_height * aspect
                
            # Center coordinates inside bounding box
            x_offset = (max_width - width) / 2
            y_offset = (max_height - height) / 2
            
            return slide.shapes.add_picture(
                img_path, 
                left + x_offset, 
                top + y_offset, 
                width=width, 
                height=height
            )
        except Exception as e:
            print(f"PIL processing error for {img_name}: {e}. Falling back to default scaling.")
            
    # Default fallback
    try:
        return slide.shapes.add_picture(img_path, left, top, width=max_width)
    except Exception as e:
        print(f"Failed to add image {img_name}: {e}")
    return None

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Common layout: blank
    blank_layout = prs.slide_layouts[6]
    
    # Theme Colors (Option B: Terracotta & Sage)
    TERRACOTTA = RGBColor(186, 85, 66)  # Warm Terracotta
    SAGE = RGBColor(110, 145, 115)        # Sage Green
    IVORY = RGBColor(253, 251, 246)       # Soft Ivory Background
    WHITE = RGBColor(255, 255, 255)
    TEXT_MAIN = RGBColor(64, 56, 50)      # Deep Browny Charcoal for high readability
    TEXT_MUTED = RGBColor(115, 105, 95)  # Warm Muted Gray
    
    # ==========================================
    # SLIDE 1: COVER SLIDE (Terracotta Background)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, TERRACOTTA)
    
    # Top Sage accent line
    top_accent = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.15)
    )
    top_accent.fill.solid()
    top_accent.fill.fore_color.rgb = SAGE
    top_accent.line.fill.background()
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "御苑福膳"
    p1.font.name = "SimHei"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = SAGE
    p1.alignment = PP_ALIGN.LEFT
    p1.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "专注本土超级食物 · 开启健康膳食新风尚"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT
    
    # Subtitle / Company Box
    company_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
    tf2 = company_box.text_frame
    tf2.word_wrap = True
    
    pc1 = tf2.paragraphs[0]
    pc1.text = "北京中卫神州生物科技有限公司 创立"
    pc1.font.name = "Microsoft YaHei"
    pc1.font.size = Pt(16)
    pc1.font.color.rgb = RGBColor(238, 220, 210)
    pc1.space_after = Pt(8)
    
    pc2 = tf2.add_paragraph()
    pc2.text = "北京御苑福膳文化传播有限公司 运营管理"
    pc2.font.name = "Microsoft YaHei"
    pc2.font.size = Pt(16)
    pc2.font.color.rgb = RGBColor(238, 220, 210)

    # ==========================================
    # SLIDE 2: BRAND ORIGIN & RESEARCH BACKGROUND (Ivory Background)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, IVORY)
    add_header(slide2, "品牌起源与中医药科研背景")
    
    # Left Box - Description Card
    left_card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(235, 230, 222)
    
    left_tf = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.4)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "深厚中医药学科研背景"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)
    
    p = left_tf.add_paragraph()
    p.text = "“御苑福膳”由北京中卫神州生物科技有限公司（成立于2002年，原隶属中国医学科学院药用植物研究所企业）创立。现由其子公司北京御苑福膳文化传播有限公司运营管理。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(10)
    
    p = left_tf.add_paragraph()
    p.text = "依托数十年中医药科研的历史积淀，跨界联手现代营养学专家，倾力构建餐桌健康工程，打造全民康养与健康膳食餐饮标杆。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MAIN
    
    # Left Card Image (Research Base)
    add_image_fit(slide2, "base_research.jpg", Inches(0.8), Inches(4.0), Inches(5.4), Inches(2.4))
    
    # Right Box - Expert Guidance List
    right_card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(235, 230, 222)
    
    right_tf = slide2.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.4)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "国医大师与权威专家学术指导"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)
    
    bullet_points = [
        "中药学专家、工程院院士肖培根长期担纲学术指导，确保底蕴；",
        "国医大师、工程院院士王琦亲临指导，对时令膳食开发提供科学支撑；",
        "汇聚多位资深中医药学者、现代营养师与烹饪大师构建健康膳食模式。"
    ]
    for bp in bullet_points:
        p = right_tf.add_paragraph()
        p.text = "• " + bp
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    # Right Card Images (Xiao Peigen & Wang Qi)
    add_image_fit(slide2, "xiao_peigen.jpg", Inches(7.1), Inches(3.9), Inches(2.5), Inches(2.1))
    add_image_fit(slide2, "wang_qi.jpg", Inches(9.9), Inches(3.9), Inches(2.5), Inches(2.1))
    
    # Right Card Image Captions
    cap_xiao = slide2.shapes.add_textbox(Inches(7.1), Inches(6.05), Inches(2.5), Inches(0.4)).text_frame
    cap_xiao.word_wrap = True
    p = cap_xiao.paragraphs[0]
    p.text = "肖培根 院士"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    
    cap_wang = slide2.shapes.add_textbox(Inches(9.9), Inches(6.05), Inches(2.5), Inches(0.4)).text_frame
    cap_wang.word_wrap = True
    p = cap_wang.paragraphs[0]
    p.text = "王琦 院士"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 3: DIETARY CONCEPT (Ivory Background)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, IVORY)
    add_header(slide3, "核心理念：本土超级食物与新鲜烹饪")
    
    # Left Column
    left_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(235, 230, 222)
    
    left_tf = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(5.0)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "新鲜食材 · 新鲜烹饪"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(16)
    
    points_left = [
        "坚持本真膳食体验，拒绝工业预制，完全遵循“新鲜食材、现场即时烹饪”的工艺准则；",
        "精选并推广高营养素密度的中国本土超级食物，将日常朴素食材以极致的烹饪艺术做出千姿百态的花样；",
        "紧密围绕二十四节气的时令更替进行菜谱轮换，锁住新鲜食材的原始生机、活性成分与营养；",
        "以健康为底色，把大众餐桌的日常饮食升级为兼具美味享受与健康调理的精致体验。"
    ]
    for pt in points_left:
        p = left_tf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(14)
        
    # Right Column
    right_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(235, 230, 222)
    
    right_tf = slide3.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "时令烹饪与健康食疗创新"
    p.font.name = "SimHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)
    
    points_right = [
        "首创“鲜品入膳”核心模式，探索生鲜本草植物与现代烹饪的有机融合；",
        "依托高标准大棚与生态基地，现采现用藿香、紫苏、玉竹等天然时令植物；",
        "科学融合中医药饮食理论与现代临床营养学成果，精心守护亚健康大众餐桌。"
    ]
    for pt in points_right:
        p = right_tf.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)
        
    # Right Card Image (Greenhouse Base)
    add_image_fit(slide3, "base_greenhouse.jpg", Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.2))
    
    # Image Caption
    cap_gh = slide3.shapes.add_textbox(Inches(7.0), Inches(6.22), Inches(5.5), Inches(0.4)).text_frame
    cap_gh.word_wrap = True
    p = cap_gh.paragraphs[0]
    p.text = "生态示范基地 · 现代化高科技大棚"
    p.font.name = "SimHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 4: SOCIAL RESPONSIBILITY & COMMUNITY (Ivory Background)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, IVORY)
    add_header(slide4, "响应大健康政策与温情回馈社会")
    
    # Left Column - Weight Management & Holiday outreach
    left_card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(235, 230, 222)
    
    left_tf = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.6)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "积极响应政策与践行社会责任"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(10)
    
    p = left_tf.add_paragraph()
    p.text = "「大健康战略与体重管理年」：长期致力于“教你吃、教你做”的科学食育与体重管理推广。每逢清明节（低卡青团）、儿童节（趣味营养体重管理膳食）、端午节（草木香草粽）、中秋节（低糖低脂养生月饼）等，创办节日美食手工工坊与推广活动，传递科学健康膳食智慧。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MAIN
    
    # Left Card Image (Academic Conference Group Photo)
    add_image_fit(slide4, "academic_conference.jpg", Inches(1.1), Inches(4.2), Inches(4.8), Inches(2.0))
    
    cap_conf = slide4.shapes.add_textbox(Inches(1.1), Inches(6.22), Inches(4.8), Inches(0.45)).text_frame
    p = cap_conf.paragraphs[0]
    p.text = "中医药膳学术研讨与品鉴会现场合影"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    
    # Right Column - Elderly Canteen
    right_card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(235, 230, 222)
    
    right_tf = slide4.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.6)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "「“爱心拐杖”云端老年食堂」保障"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(10)
    
    p = right_tf.add_paragraph()
    p.text = "作为温情回馈社会的责任企业，御苑福膳非常荣幸能够成为海淀区马连洼街道“爱心拐杖”云端老年食堂首批定点保障单位（该项目曾获《中国日报》等主流媒体专题报道与点赞）。我们以此为支点，持续深入开展健康饮食进社区活动，为长者研发营养合理、新鲜烹饪的适老膳食。"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_MAIN
    
    # Right Card Image (Canteen meeting / cooperation photo)
    add_image_fit(slide4, "representing_conference.jpg", Inches(7.3), Inches(4.2), Inches(4.8), Inches(2.0))
    
    cap_rep = slide4.shapes.add_textbox(Inches(7.3), Inches(6.22), Inches(4.8), Inches(0.45)).text_frame
    p = cap_rep.paragraphs[0]
    p.text = "社会公益保障与行业大会工作交流留影"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 5: HONORS & QUALIFICATIONS (Ivory Background)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, IVORY)
    add_header(slide5, "品牌资质与社会荣誉认可")
    
    # 5 Honors layout - 3 on top, 2 centered bottom
    honor_w = Inches(3.8)
    honor_h = Inches(2.6)
    gap_x = Inches(0.3)
    
    honors_list = [
        ("非物质文化遗产", "ich_certificate.jpg", "传统膳食技艺非遗传承保护单位", "多年来始终坚持中华古老膳食技艺的精髓继承与现代保护，非物质文化遗产的传承，见证了御苑福膳深厚的文化底蕴。"),
        ("老字号荣誉", "time_honored_brand.jpg", "传统膳食品牌与老字号荣誉", "荣获中医药及健康餐饮行业内的老字号名号，代表着数十载持之以恒的卓越声誉、工匠精神以及官方与社会的信任。"),
        ("行业荣誉资质", "honor_plaque.jpg", "健康膳食与特色康养餐饮认证", "被行业协会与健康养生协会联合认证为健康膳食与特色健康餐饮示范店，体现出深厚的草本膳食科学研究底蕴。"),
        ("首都窗口行业示范单位", "skills_demo.png", "2008年北京迎奥运技能示范", "在北京市劳动和社会保障局、首都窗口行业奥运培训工作协调小组联合评比中，被授予“首都窗口行业技能示范单位”牌匾。"),
        ("安全生产标准化企业", "safety_production.png", "餐饮业安全生产标准化企业", "作为北京市海淀区餐饮业安全生产标准化的典范，连续多年荣获安全生产达标资质，用极高安全标准保障每一位食客的健康安全。")
    ]
    
    for i, (title, img_name, subtitle, desc) in enumerate(honors_list):
        if i < 3:
            x = Inches(0.6) + i * (honor_w + gap_x)
            y = Inches(1.4)
        else:
            # Centered on bottom row
            x = Inches(2.65) + (i - 3) * (honor_w + gap_x)
            y = Inches(4.2)
            
        # Card background
        card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, honor_w, honor_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(235, 230, 222)
        
        # Embedded certificate/plaque image
        add_image_fit(slide5, img_name, x + Inches(0.15), y + Inches(0.25), Inches(1.3), Inches(2.1))
        
        # Text Frame
        tf = slide5.shapes.add_textbox(x + Inches(1.5), y + Inches(0.15), Inches(2.2), Inches(2.3)).text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "SimHei"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TERRACOTTA
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = SAGE
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: CELEBRITY FOOTPRINTS (Ivory Background)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, IVORY)
    add_header(slide6, "社会各界与中医药学专家交流足迹")
    
    # Left Box - Political leaders
    left_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = RGBColor(235, 230, 222)
    
    left_tf = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(2.6)).text_frame
    left_tf.word_wrap = True
    
    p = left_tf.paragraphs[0]
    p.text = "党政领导与社会名流莅临视察"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)
    
    leaders = [
        "全国人大常委会副委员长 何维 莅临考察交流健康膳食并亲切留影；",
        "全国政协原副主席 张思卿 & 孙孚凌 视察指导工作；",
        "原卫生部部长 张文康 & 高强、原副部长 顾英奇 亲临指导；",
        "甘肃省委原书记 陆浩 莅临视察指导健康膳食工作并亲切交流；",
        "海南省委原书记 卫留成 等先后莅临考察；",
        "毛泽东主席原机要秘书 张玉凤 女士及刘爱民先生亲临品鉴。"
    ]
    for ld in leaders:
        p = left_tf.add_paragraph()
        p.text = "★ " + ld
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(5)
        
    # Left Column Photos (He Wei & Lu Hao)
    add_image_fit(slide6, "he_wei.jpg", Inches(0.8), Inches(4.2), Inches(2.6), Inches(2.0))
    add_image_fit(slide6, "lu_hao.jpg", Inches(3.6), Inches(4.2), Inches(2.6), Inches(2.0))
    
    cap_he = slide6.shapes.add_textbox(Inches(0.8), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_he.paragraphs[0]
    p.text = "何维 副委员长莅临考察"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    
    cap_lu = slide6.shapes.add_textbox(Inches(3.6), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_lu.paragraphs[0]
    p.text = "原省委书记 陆浩视察留影"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
        
    # Right Box - Experts & Masters
    right_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = RGBColor(235, 230, 222)
    
    right_tf = slide6.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.6)).text_frame
    right_tf.word_wrap = True
    
    p = right_tf.paragraphs[0]
    p.text = "学术泰斗与国宴大师技术合作"
    p.font.name = "SimHei"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)
    
    experts = [
        "中药学专家 肖培根院士 亲自指导，从本草药理角度确保膳食调理；",
        "国医大师 王琦院士 亲临指导健康饮食与时令膳食应用；",
        "中国工程院原副院长 刘德培院士 亲临指导餐桌大健康建设；",
        "著名病毒学家 顾方舟院士 莅临并对大众健康餐饮模式予以高度肯定；",
        "国宴烹饪大师 侯仲华 & 孟宪斌大师 指导，将传统草本与国宴技艺结合。"
    ]
    for exp in experts:
        p = right_tf.add_paragraph()
        p.text = "★ " + exp
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(5)
        
    # Right Column Photos (Liu Depei & Hou Zhonghua)
    add_image_fit(slide6, "liu_depei.jpg", Inches(7.0), Inches(4.2), Inches(2.6), Inches(2.0))
    add_image_fit(slide6, "hou_zhonghua.jpg", Inches(9.8), Inches(4.2), Inches(2.6), Inches(2.0))
    
    cap_liu = slide6.shapes.add_textbox(Inches(7.0), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_liu.paragraphs[0]
    p.text = "刘德培 院士亲临指导"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    
    cap_hou = slide6.shapes.add_textbox(Inches(9.8), Inches(6.22), Inches(2.6), Inches(0.45)).text_frame
    p = cap_hou.paragraphs[0]
    p.text = "国宴烹饪大师 侯仲华指导"
    p.font.name = "SimHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 7: END SLIDE (Terracotta Background)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, TERRACOTTA)
    
    # Bottom Sage accent line
    bottom_accent = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35), Inches(13.333), Inches(0.15)
    )
    bottom_accent.fill.solid()
    bottom_accent.fill.fore_color.rgb = SAGE
    bottom_accent.line.fill.background()
    
    # Text Box
    end_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf_end = end_box.text_frame
    tf_end.word_wrap = True
    
    p = tf_end.paragraphs[0]
    p.text = "坚持新鲜食材与新鲜烹饪"
    p.font.name = "SimHei"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)
    
    p = tf_end.add_paragraph()
    p.text = "以时令膳食守护大众餐桌健康"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(40)
    
    p = tf_end.add_paragraph()
    p.text = "谢谢观看"
    p.font.name = "SimHei"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SAGE
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filepath = os.path.join(proj_dir, "御苑福膳_品牌介绍_精简版.pptx")
    prs.save(filepath)
    print(f"Presentation generated successfully at {filepath}")

if __name__ == "__main__":
    create_presentation()
